"""
presentation.py - Module for PowerPoint Presentation Creation.

This module provides functions to create and modify PowerPoint presentations based on provided templates,
user input, and language model responses.

Module Functions:
    - find_most_similar_layout(prs, target_name):
        Find the most similar layout in the PowerPoint presentation to the target layout name.

    - find_content_placeholder(slide):
        Find the content placeholder in a slide.

    - find_slide_layout_by_name(prs, layout_name):
        Find a slide layout by its name in the PowerPoint presentation.

    - create_presentation(template_file, slide_content, output_file, logger):
        Create a new PowerPoint presentation based on a template and slide content.

    - build_slides_with_llm(template_file, word_content, output_file, session_file, logger):
        Build slides using a language model (LLM) based on user input.

    - get_latest_slide_deck(pptx_file_path, current_slide_deck, logger):
        Retrieves information from the latest PowerPoint slide deck in the specified file.


Dependencies:
    - difflib
    - json
    - os
    - pickle
    - pptx (from python-pptx package)
    - pptx.enum.shapes.MSO_SHAPE_TYPE
    - utils (custom utility functions in utils.py)
    - llm_ops (custom functions in llm_ops.py)

"""

import difflib
import json
import os
import pickle
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .utils import format_error, format_info, format_prompt, format_warning, ensure_list
from .llm_ops import get_chat_response




def find_most_similar_layout(prs, target_name):
    """Find the most similar layout in the PowerPoint presentation to the target layout name.

    Args:
        prs (pptx.Presentation): The PowerPoint presentation object.
        target_name (str): The name of the layout to find.

    Returns:
        pptx.slide.SlideLayout or None: The most similar layout if found, else None.
    """
    layout_names = [layout.name for layout in prs.slide_layouts]
    closest_matches = difflib.get_close_matches(target_name, layout_names)

    if closest_matches:
        closest_match = closest_matches[0]
        for layout in prs.slide_layouts:
            if layout.name == closest_match:
                return layout

    return None

def find_content_placeholder(slide):
    """Find the content placeholder in a slide.

    Args:
        slide (pptx.slide.Slide): The slide to search for the content placeholder.

    Returns:
        pptx.shapes.placeholder.Placeholder or None: The content placeholder if found, else None.
    """
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.placeholder_format.idx == 1:
            return shape
    return None


def find_slide_layout_by_name(prs, layout_name):
    """Find a slide layout by its name in the PowerPoint presentation.

    Args:
        prs (pptx.Presentation): The PowerPoint presentation object.
        layout_name (str): The name of the layout to find.

    Returns:
        pptx.slide.SlideLayout or None: The slide layout if found, else None.
    """
    for slide_layout in prs.slide_layouts:
        if slide_layout.name == layout_name:
            return slide_layout
    return None


def create_presentation(template_file, slide_content, output_file,logger):
    """Create a new PowerPoint presentation based on a template and slide content.

    Args:
        template_file (str): The path to the pptx PowerPoint file to use the template from.
        slide_content (list): A list of dictionaries representing slide content.
        output_file (str): The path to save the new PowerPoint presentation.
        logger (logging.Logger): The logger object to record any errors.
    """
    # Load the template PowerPoint file
    prs = Presentation(template_file)

    # Delete all slides from the presentation
    slides_to_remove = prs.slides._sldIdLst[:]
    slide_ids_to_remove = [slide.rId for slide in prs.slides._sldIdLst]
    for slide_id in slide_ids_to_remove:
        prs.part.drop_rel(slide_id)
    for slide_id in slides_to_remove:
        prs.slides._sldIdLst.remove(slide_id)

    title_and_content_layout = find_most_similar_layout(prs, "Title and Content")
    if not title_and_content_layout:
        print(format_warning("Warning: Layout 'Title and Content' not found in the template."))
        return

    for slide_data in slide_content:
        # Create a new slide using the original layout from the template
        slide = prs.slides.add_slide(title_and_content_layout)

        # Set the title and content of the new slide
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_data.get("title", "")

        content_placeholder = find_content_placeholder(slide)
        if content_placeholder:
            content_placeholder.text = slide_data.get("content", "")

        # Add presentation transcript to the notes section of each slide
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data.get("narration", "")


    # Save the slides gracefully
    while True:
        try:
            prs.save(output_file)
            break
        except PermissionError:
            print(format_warning("Action Required:")+ f" The file {output_file} is open in PowerPoint or another application." +format_prompt("Please close it and press Enter to retry. "))
            logger.error(f"Error: The file '{output_file}' is open in PowerPoint or another application. Please close it and press Enter to retry.")
            input()  # Wait for the user to close the file
        except Exception as e:
            print(format_error("An unexpected error occurred while saving the presentation: ")+f"{e}")
            logger.error(f"An unexpected error occurred while saving the presentation: {e}")
            break
    print(format_info("Presentation created successfully and saved as: ")+ f"{output_file}")
    # Open the saved presentation file
    try:
        os.startfile(output_file)
    except Exception as e:
            print(format_info("Open the saved file manually to see the changes."))
            logger.debug(f"Can not preview the presentation using os default viewer: {e}")


def get_latest_slide_deck(pptx_file_path, current_slide_deck, logger):
    """
    Retrieves information from the latest PowerPoint slide deck in the specified file.

    Args:
        pptx_file_path (str): The path to the PowerPoint file.
        current_slide_deck (list): The current list of slide information.
        logger (Logger): The logger object for error reporting.

    Returns:
        list: A list of dictionaries containing slide information, including slide number,
              title, content, and narration.

    Note:
        This function attempts to load the PowerPoint presentation and extract slide
        information. If the file is open in another application with a read lock (rare), it will prompt the user
        to close it and retry. If the file does not exist, it returns the current slide deck.

    """
    presentation = None
    while True:
        try:
            presentation = Presentation(pptx_file_path)
            break
        except PermissionError:
            print(format_warning("Action Required:")+ f" The file {output_file} is open in PowerPoint or another application." +format_prompt("Please close it and press Enter to retry. "))
            logger.error(f"Error: The file '{output_file}' is open in PowerPoint or another application. Please close it and press Enter to retry.")
            input()  # Wait for the user to close the file
        except FileNotFoundError:
            return current_slide_deck           
        except Exception as e:
            logger.error(f"An unexpected error occurred while loading the presentation: {e}")
            return current_slide_deck            

    updated_slide_deck = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_info = {"slide_number": float(slide_number), "title": "", "content": "", "narration": ""}
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape == slide.shapes[0]:  # Assuming the first shape is the title
                    slide_info["title"] = shape.text
                else:
                    slide_info["content"] += shape.text + "\n"
        # Extract notes from the notes slide
        notes_slide = slide.notes_slide
        for shape in notes_slide.shapes:
            if shape.has_text_frame:
                slide_info["narration"] += shape.text + "\n"

        updated_slide_deck.append(slide_info)

    return updated_slide_deck


def build_slides_with_llm(template_file,word_content, output_file, session_file, logger):
    """Build slides using a language model (LLM) based on user input.

    Args:
        template_file (str): The path to the template PowerPoint file.
        word_content (str): The content for the language model input.
        output_file (str): The path to save the final PowerPoint presentation.
        session_file (str): The path to the session file to resume a previously saved session.
        logger (logging.Logger): The logger object to record any errors.
    """
    UIM=["\nHave a look at what we have got so far. \nKey in -1 to undo, or let me know if you want me to make further changes: "]
    slide_deck_history=[[{"slide_number": 1, "title": "", "content": ""}]]
    try:
        if not session_file:
            initial_slide_deck = get_latest_slide_deck(output_file, slide_deck_history[0], logger)
            if initial_slide_deck != slide_deck_history[0]:
                print(format_warning("The output file is not empty but no session file is supplied. If you have the session file, consider continuing from the previously saved session.")+"\n")
                print(format_info("Loading the existing text content of the file into a new session...")+"\n")
                slide_deck = initial_slide_deck
                slide_deck_history=[slide_deck]
            else:
                # Ask the user for prompt
                prompt = input(format_prompt("\nWhat shall I do for you? "))
                print(format_info("OK, please wait. This might take a while...")+"\n")
                slide_deck=get_chat_response(word_content, initial_slide_deck, prompt,logger)
                logger.debug(slide_deck)
                while True:
                    try:
                        slide_deck=ensure_list(json.loads(slide_deck))
                        slide_deck_history=[slide_deck]
                        break
                    except Exception as e:
                        logger.debug("Not neccessarily an error but Invalid JSON string")
                        logger.debug(e)
                        # Ask the user for prompt
                        prompt = input(format_prompt("Hmm, not sure what you want so I did not make any changes. Try differently: "))
                        print(format_info("OK, please wait. This might take a while...")+"\n")
                        slide_deck=get_chat_response(word_content,[], prompt,logger)                    
                    


        else:
            print(format_info("Resuming a previously saved session ...")+"\n")
            logger.debug("Resuming a previously saved session.")
            # Load session data from the session file
            with open(session_file, 'rb') as file:
                loaded_data = pickle.load(file)
                logger.debug(loaded_data)
            # Access slide_deck_history variable from the loaded data
            slide_deck_history = loaded_data['slide_deck_history']
            slide_deck = get_latest_slide_deck(output_file, slide_deck_history[-1], logger)
            slide_deck_history[-1] = slide_deck
            logger.debug(slide_deck)

        while True:
            create_presentation(template_file,slide_deck,output_file,logger)
            prompt = input(format_prompt(UIM.pop() if len(UIM) > 1 else UIM[-1]))
            print(format_info("OK, please wait. This might take a while...")+"\n")
            if prompt == "-1":
                print(format_info("Undoing changes and loading earlier version ...")+"\n")
                slide_deck_history.pop() if len(slide_deck_history) > 1 else None
                slide_deck = slide_deck_history[-1]
                continue
            slide_deck = get_latest_slide_deck(output_file, slide_deck, logger)
            result=get_chat_response(word_content,slide_deck, prompt,logger)
            try:
                result=ensure_list(json.loads(result))
            except Exception as e:
                logger.debug("Not neccessarily an error but Invalid JSON string")
                logger.debug(e)
                UIM.append("\nHmm. Not sure what you want so I did not make any changes. Try differently: ")
                continue

            logger.debug("******* This is the result:*********\n")
            logger.debug(result)
            # Convert slide numbers to strings in both arrays
            for slide in slide_deck:
                slide["slide_number"] = "{:.1f}".format(float(slide["slide_number"]))

            for slide in result:
                slide["slide_number"] = "{:.1f}".format(float(slide["slide_number"]))

            # Remove existing slides with duplicate slide numbers from slide_deck
            slide_numbers_in_new_slides = set(slide["slide_number"] for slide in result)
            slide_deck = [slide for slide in slide_deck if slide["slide_number"] not in slide_numbers_in_new_slides]

            # Create sets of slide numbers and their opposites for both arrays
            slide_numbers_in_slide_deck = set(slide["slide_number"] for slide in slide_deck)
            opposite_slide_numbers_in_slide_deck = set("{:.1f}".format(-float(num)) for num in slide_numbers_in_slide_deck)
            logger.debug(opposite_slide_numbers_in_slide_deck)

            slide_numbers_in_result = set(slide["slide_number"] for slide in result)
            opposite_slide_numbers_in_result = set("{:.1f}".format(-float(num)) for num in slide_numbers_in_result)
            logger.debug(opposite_slide_numbers_in_result)
            # Remove slides with slide numbers matching opposites from both arrays
            filtered_slide_deck = [slide for slide in slide_deck if slide["slide_number"] not in opposite_slide_numbers_in_result]
            filtered_result = [slide for slide in result if slide["slide_number"] not in opposite_slide_numbers_in_slide_deck]

            # Merge the two arrays while preserving order
            merged_slides = sorted(filtered_slide_deck + filtered_result, key=lambda x: float(x["slide_number"]))

            # Rewrite the slide numbers starting from 1
            for i, slide in enumerate(merged_slides, start=1):
                slide["slide_number"] = i

            slide_deck = merged_slides
            slide_deck_history.append(slide_deck)
            logger.debug("******* This is the slide_deck:*********\n")
            logger.debug(slide_deck)
    except Exception as e:
        logger.error("Something went wrong while making presentation:")
        logger.error(e)
        print(format_error("Something went wrong while making presentation: ")+f"{e}"+"\n")
    finally:
        file_name=output_file.replace(".", "_")+'_session.pkl'
        if os.path.exists(file_name) and not session_file:
            file_name = file_name.replace(".pkl", f"_new.pkl")
        print("\nSaving the session to "+file_name)
        # Saving variables to a file
        with open(file_name, 'wb') as file:
            session_data = {
                'slide_deck_history': slide_deck_history,
                'word_content': word_content
            }
            pickle.dump(session_data, file)
            # Flush the buffer to ensure data is written to the file immediately
            file.flush()
